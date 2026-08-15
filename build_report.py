from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEAL_DARK = RGBColor(0x0B, 0x3D, 0x3A)
GOLD = RGBColor(0xD4, 0xA6, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# ============================================================
# SLIDE 1 : TITRE
# ============================================================
slide = prs.slides.add_slide(blank_layout)

# Fond teal fonce
bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)  # 1 = rectangle
bg.fill.solid()
bg.fill.fore_color.rgb = TEAL_DARK
bg.line.fill.background()
bg.shadow.inherit = False

# Sur-titre
box = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(10), Inches(0.4))
tf = box.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "TOGO AI LAB - DATA CHALLENGE ENVIRONNEMENT - DEFI 1"
run.font.size = Pt(13)
run.font.bold = True
run.font.color.rgb = GOLD

# Titre principal
box = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11), Inches(1.6))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Diagnostic de l'acces a l'eau potable au Togo"
run.font.size = Pt(38)
run.font.bold = True
run.font.color.rgb = WHITE

prs.save("rapport_defi_eau.pptx")
print("PowerPoint sauvegarde : rapport_defi_eau.pptx")